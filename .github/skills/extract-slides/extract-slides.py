#!/usr/bin/env python3
"""Extract text content from each slide in the PowerPoint file and write
individual Markdown files to the text-extract directory.

Produces structured Markdown with proper headings for slide titles,
bullet-point lists respecting indentation levels, and Markdown tables
for any table shapes.

Usage:
    python .github/skills/extract-slides/extract-slides.py [path/to/file.pptx]

If no path is provided, defaults to 'content/Kubernetes Technical Primer.pptx'
(resolved relative to the repository root).
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is not installed. Install it with:\n  pip install python-pptx")
    sys.exit(1)


# Placeholder type indices (from the Open XML spec / python-pptx)
# Used to identify the role of each placeholder shape on a slide.
_PH_TITLE = 0       # Title
_PH_CENTER_TITLE = 3  # Center title (used on title slides)
_PH_SUBTITLE = 1     # Subtitle / body on title slides
_PH_BODY = 2         # Body / content
_PH_TABLE = 12       # Table placeholder


def _placeholder_idx(shape):
    """Return the placeholder index for a shape, or None if not a placeholder."""
    try:
        return shape.placeholder_format.idx
    except (AttributeError, ValueError):
        return None


def _is_title_shape(shape):
    """Return True if the shape is a title or center-title placeholder."""
    idx = _placeholder_idx(shape)
    if idx in (_PH_TITLE, _PH_CENTER_TITLE):
        return True
    # Fallback: check the shape name (works for non-standard layouts)
    name = (shape.name or "").lower()
    return "title" in name and "subtitle" not in name


def _is_subtitle_shape(shape):
    """Return True if the shape is a subtitle placeholder."""
    idx = _placeholder_idx(shape)
    if idx == _PH_SUBTITLE:
        return True
    name = (shape.name or "").lower()
    return "subtitle" in name


def _collapse_text(text):
    """Collapse whitespace in a string, stripping leading/trailing space."""
    return " ".join(text.split())


def _extract_text_frame(shape, heading_level=None):
    """Extract paragraphs from a shape's text frame as Markdown lines.

    - If heading_level is given (e.g. 2 or 3), the first non-empty paragraph
      is emitted as a Markdown heading at that level, and subsequent
      paragraphs are emitted as body text.
    - Bullet / indentation levels (paragraph.level) are rendered as nested
      Markdown unordered list items.
    - Paragraphs with no indentation level are rendered as plain text.
    """
    if not shape.has_text_frame:
        return []

    lines = []
    heading_emitted = heading_level is None  # skip heading logic if not requested

    for para in shape.text_frame.paragraphs:
        text = _collapse_text(para.text)
        if not text:
            continue

        # Emit the first paragraph as a heading if requested
        if not heading_emitted:
            lines.append(f"{'#' * heading_level} {text}")
            heading_emitted = True
            continue

        level = para.level or 0

        # Detect if the paragraph looks like a bullet / list item.
        # python-pptx exposes paragraph.level (0–8) for indented content.
        # Level 0 with no bullet can be plain text; level >= 1 is always a
        # nested list item.  We treat level-0 items inside a body/content
        # placeholder as top-level bullets if there are multiple paragraphs.
        if heading_level is not None:
            # Inside a titled shape — lines after the heading are bullets
            indent = "  " * level
            lines.append(f"{indent}- {text}")
        elif level > 0:
            indent = "  " * (level - 1)
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)

    return lines


def _extract_table(shape):
    """Extract a table shape as a Markdown table."""
    if not shape.has_table:
        return []

    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(_collapse_text(cell.text))
        rows.append(cells)

    if not rows:
        return []

    lines = []
    # Header row
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    # Data rows
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return lines


def _iter_shapes(shapes):
    """Yield all shapes, recursing into grouped shapes."""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_shapes(shape.shapes)
            except AttributeError:
                pass
        else:
            yield shape


def extract_slide_markdown(slide, slide_number):
    """Return the full Markdown content for a single slide.

    Output structure:
      # <Title>             (if the slide has a title placeholder)
      ## <Subtitle>         (if present)
      <body content with bullets>
      <tables>
      <other text shapes>
    """
    sections = []  # list of (priority, lines) for ordering

    title_shape = None
    subtitle_shape = None
    body_shapes = []
    table_shapes = []
    other_shapes = []

    for shape in _iter_shapes(slide.shapes):
        if _is_title_shape(shape):
            title_shape = shape
        elif _is_subtitle_shape(shape):
            subtitle_shape = shape
        elif shape.has_table:
            table_shapes.append(shape)
        elif shape.has_text_frame:
            # Distinguish body/content placeholders from free-floating text boxes
            idx = _placeholder_idx(shape)
            if idx is not None and idx not in (_PH_TITLE, _PH_CENTER_TITLE, _PH_SUBTITLE):
                body_shapes.append(shape)
            else:
                other_shapes.append(shape)
        # Images and other non-text shapes are skipped

    md_lines = []

    # --- Title ---
    if title_shape and title_shape.has_text_frame:
        title_text = _collapse_text(title_shape.text_frame.text)
        if title_text:
            md_lines.append(f"# {title_text}")

    # --- Subtitle ---
    if subtitle_shape and subtitle_shape.has_text_frame:
        subtitle_text = _collapse_text(subtitle_shape.text_frame.text)
        if subtitle_text:
            if md_lines:
                md_lines.append("")
            md_lines.append(f"## {subtitle_text}")

    # --- Body / content placeholders ---
    for shape in body_shapes:
        lines = _extract_text_frame(shape)
        if lines:
            md_lines.append("")
            # If all lines are plain text (no bullets), check paragraph levels
            # to decide whether to render as a bullet list
            has_multi = len(lines) > 1
            for para in shape.text_frame.paragraphs:
                text = _collapse_text(para.text)
                if not text:
                    continue
                level = para.level or 0
                indent = "  " * level
                if has_multi:
                    md_lines.append(f"{indent}- {text}")
                else:
                    md_lines.append(text)

    # --- Tables ---
    for shape in table_shapes:
        table_lines = _extract_table(shape)
        if table_lines:
            md_lines.append("")
            md_lines.extend(table_lines)

    # --- Other text shapes (free text boxes, callouts, etc.) ---
    for shape in other_shapes:
        lines = _extract_text_frame(shape)
        if lines:
            md_lines.append("")
            md_lines.extend(lines)

    md_lines.append("")  # trailing newline
    return "\n".join(md_lines)


def main():
    # Resolve repo root (three levels up: .github/skills/extract-slides/ -> repo root)
    script_dir = Path(__file__).resolve().parent
    repo_root = script_dir.parent.parent.parent

    # Determine PPTX path
    if len(sys.argv) > 1:
        pptx_path = Path(sys.argv[1])
    else:
        pptx_path = repo_root / "content" / "Kubernetes Technical Primer.pptx"

    if not pptx_path.exists():
        print(f"ERROR: PowerPoint file not found: {pptx_path}")
        sys.exit(1)

    output_dir = repo_root / "content" / "text-extract"
    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Reading: {pptx_path}")

    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)
    print(f"Found {slide_count} slides")

    # Remove any old slide-*.md files that may be stale (e.g. slide count decreased)
    for old_file in output_dir.glob("slide-*.md"):
        old_file.unlink()

    for idx, slide in enumerate(prs.slides, start=1):
        content = extract_slide_markdown(slide, idx)

        out_path = output_dir / f"slide-{idx}.md"
        out_path.write_text(content, encoding="utf-8")

    print(f"Wrote {slide_count} files to content/text-extract/")


if __name__ == "__main__":
    main()
